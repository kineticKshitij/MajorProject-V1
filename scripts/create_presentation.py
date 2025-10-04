#!/usr/bin/env python
"""
Generate a comprehensive PowerPoint presentation for the AI-Enhanced Google Dorks Toolkit
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    """Create a comprehensive PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    
    # Define color scheme
    PRIMARY_COLOR = RGBColor(0, 123, 255)    # Blue
    SECONDARY_COLOR = RGBColor(108, 117, 125) # Gray
    SUCCESS_COLOR = RGBColor(40, 167, 69)     # Green
    WARNING_COLOR = RGBColor(255, 193, 7)    # Yellow
    DANGER_COLOR = RGBColor(220, 53, 69)     # Red
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "AI-Enhanced Google Dorks Toolkit"
    subtitle.text = "A Comprehensive Cybersecurity Research Platform\n\nDeveloped by: [Your Name]\nDate: October 2025"
    
    # Format title
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.bold = True
    
    # Slide 2: Project Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Project Overview"
    
    content = slide2.placeholders[1].text_frame
    content.text = "🎯 Mission: Democratize cybersecurity research through intelligent Google dorking"
    
    p = content.add_paragraph()
    p.text = "🔍 What are Google Dorks?"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "Specialized search queries that reveal sensitive information exposed on the internet"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "🤖 AI Integration"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "Google Gemini-powered chatbot for intelligent security guidance"
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "🛡️ Ethical Focus"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "Educational tool for security professionals and researchers"
    p.level = 2
    
    # Slide 3: Key Features
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "🚀 Key Features"
    
    content = slide3.placeholders[1].text_frame
    content.text = "📚 Comprehensive Dork Library"
    
    features = [
        ("Pre-loaded with 20+ Google dorks across 8 categories", 1),
        ("Advanced search and filtering capabilities", 1),
        ("🔍 Smart Search Engine", 0),
        ("Filter by category, difficulty, risk level, and keywords", 1),
        ("Real-time search execution with Google integration", 1),
        ("🤖 AI-Powered Assistant", 0),
        ("Google Gemini integration for intelligent responses", 1),
        ("Context-aware security research guidance", 1),
        ("Personalized dork recommendations", 1),
        ("👤 User Management", 0),
        ("Custom user profiles with personal API keys", 1),
        ("Bookmark system for favorite dorks", 1),
        ("Session management for organized research", 1)
    ]
    
    for feature_text, level in features:
        p = content.add_paragraph()
        p.text = feature_text
        p.level = level
    
    # Slide 4: Technical Architecture
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "🏗️ Technical Architecture"
    
    content = slide4.placeholders[1].text_frame
    content.text = "Frontend Layer"
    
    arch_components = [
        ("🌐 Frontend Layer", 0),
        ("Bootstrap 5 responsive design", 1),
        ("AJAX for dynamic interactions", 1),
        ("Modern CSS with custom styling", 1),
        ("⚙️ Backend Layer", 0),
        ("Django 5.2.6 web framework", 1),
        ("Custom User authentication system", 1),
        ("RESTful API endpoints", 1),
        ("🗄️ Data Layer", 0),
        ("SQLite/PostgreSQL database", 1),
        ("Optimized database relationships", 1),
        ("Comprehensive data models", 1),
        ("🔌 External Integrations", 0),
        ("Google Gemini API for AI capabilities", 1),
        ("Google Search for dork execution", 1),
        ("Secure API key management", 1)
    ]
    
    for component, level in arch_components:
        p = content.add_paragraph()
        p.text = component
        p.level = level
    
    # Slide 5: Database Models
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "📊 Database Architecture"
    
    content = slide5.placeholders[1].text_frame
    content.text = "Core Models"
    
    models = [
        ("👤 User Management", 0),
        ("Custom User model with API key storage", 1),
        ("UserProfile for extended information", 1),
        ("🔍 Dork Management", 0),
        ("GoogleDork: Individual dork queries with metadata", 1),
        ("DorkCategory: Organizational categories", 1),
        ("SearchResult: Tracked search outcomes", 1),
        ("📝 Session Management", 0),
        ("SearchSession: Grouped research sessions", 1),
        ("DorkBookmark: User-saved favorites", 1),
        ("🤖 AI Integration", 0),
        ("ChatSession: AI conversation sessions", 1),
        ("ChatMessage: Individual chat messages", 1),
        ("ChatFeedback: User feedback on AI responses", 1)
    ]
    
    for model, level in models:
        p = content.add_paragraph()
        p.text = model
        p.level = level
    
    # Slide 6: User Workflows
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "🔄 User Workflows"
    
    content = slide6.placeholders[1].text_frame
    content.text = "1. Registration & Authentication"
    
    workflows = [
        ("User registration with email validation", 1),
        ("Personal API key setup for AI features", 1),
        ("Profile management and customization", 1),
        ("2. Dork Exploration", 0),
        ("Browse dorks by category or search", 1),
        ("View detailed dork information", 1),
        ("Execute dorks directly in Google", 1),
        ("Bookmark favorites for quick access", 1),
        ("3. Session Management", 0),
        ("Create organized research sessions", 1),
        ("Track progress and results", 1),
        ("Export findings in multiple formats", 1),
        ("4. AI Assistance", 0),
        ("Interactive chat with security expert AI", 1),
        ("Get personalized dork recommendations", 1),
        ("Learn ethical hacking best practices", 1)
    ]
    
    for workflow, level in workflows:
        p = content.add_paragraph()
        p.text = workflow
        p.level = level
    
    # Slide 7: Security Features
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "🛡️ Security Implementation"
    
    content = slide7.placeholders[1].text_frame
    content.text = "Authentication & Authorization"
    
    security = [
        ("🔐 User Authentication", 0),
        ("Django's built-in authentication system", 1),
        ("Session-based security", 1),
        ("Password strength requirements", 1),
        ("🔑 API Key Management", 0),
        ("Encrypted storage of personal API keys", 1),
        ("Secure key validation", 1),
        ("Individual user API isolation", 1),
        ("🛡️ Data Protection", 0),
        ("CSRF protection on all forms", 1),
        ("XSS prevention measures", 1),
        ("SQL injection protection", 1),
        ("Input validation and sanitization", 1),
        ("📋 Audit & Logging", 0),
        ("Comprehensive activity logging", 1),
        ("Security event monitoring", 1),
        ("Rate limiting implementation", 1)
    ]
    
    for sec_feature, level in security:
        p = content.add_paragraph()
        p.text = sec_feature
        p.level = level
    
    # Slide 8: AI Integration Details
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "🤖 AI Integration: Google Gemini"
    
    content = slide8.placeholders[1].text_frame
    content.text = "Intelligent Chatbot Features"
    
    ai_features = [
        ("🧠 Context-Aware Responses", 0),
        ("Understands cybersecurity terminology", 1),
        ("Provides relevant dork recommendations", 1),
        ("Maintains conversation context", 1),
        ("🎯 Specialized Knowledge", 0),
        ("Google dorking techniques and best practices", 1),
        ("Ethical hacking guidelines", 1),
        ("Security research methodologies", 1),
        ("📚 Dynamic Learning", 0),
        ("Learns from user interactions", 1),
        ("Improves responses based on feedback", 1),
        ("Adapts to user expertise level", 1),
        ("🔒 Privacy & Security", 0),
        ("User-specific API key management", 1),
        ("Secure conversation storage", 1),
        ("Data isolation between users", 1)
    ]
    
    for ai_feature, level in ai_features:
        p = content.add_paragraph()
        p.text = ai_feature
        p.level = level
    
    # Slide 9: Technology Stack
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "💻 Technology Stack"
    
    content = slide9.placeholders[1].text_frame
    content.text = "Backend Technologies"
    
    tech_stack = [
        ("🐍 Backend", 0),
        ("Python 3.13", 1),
        ("Django 5.2.6", 1),
        ("SQLite/PostgreSQL", 1),
        ("🎨 Frontend", 0),
        ("HTML5 & CSS3", 1),
        ("Bootstrap 5", 1),
        ("JavaScript (ES6+)", 1),
        ("AJAX for dynamic content", 1),
        ("🤖 AI Integration", 0),
        ("Google Gemini API", 1),
        ("Python google-genai package", 1),
        ("🛠️ Development Tools", 0),
        ("Git version control", 1),
        ("Virtual environment (venv)", 1),
        ("Django admin interface", 1),
        ("Custom management commands", 1)
    ]
    
    for tech, level in tech_stack:
        p = content.add_paragraph()
        p.text = tech
        p.level = level
    
    # Slide 10: Project Statistics
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "📈 Project Metrics"
    
    content = slide10.placeholders[1].text_frame
    content.text = "Codebase Statistics"
    
    stats = [
        ("📁 Project Structure", 0),
        ("3 Django applications (accounts, googledorks, chatbot)", 1),
        ("20+ pre-loaded Google dorks", 1),
        ("8 dork categories", 1),
        ("50+ HTML templates", 1),
        ("🔧 Code Quality", 0),
        ("Comprehensive error handling", 1),
        ("Django best practices implementation", 1),
        ("Responsive design principles", 1),
        ("Security-first approach", 1),
        ("📚 Documentation", 0),
        ("Complete workflow diagrams", 1),
        ("API documentation", 1),
        ("Setup and deployment guides", 1),
        ("Comprehensive testing suite", 1)
    ]
    
    for stat, level in stats:
        p = content.add_paragraph()
        p.text = stat
        p.level = level
    
    # Slide 11: Use Cases
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "🎯 Use Cases & Applications"
    
    content = slide11.placeholders[1].text_frame
    content.text = "Educational Applications"
    
    use_cases = [
        ("🎓 Educational Sector", 0),
        ("Cybersecurity training programs", 1),
        ("University research projects", 1),
        ("Security awareness workshops", 1),
        ("🏢 Professional Applications", 0),
        ("Penetration testing preparation", 1),
        ("Security auditing research", 1),
        ("Vulnerability assessment training", 1),
        ("👨‍💻 Research & Development", 0),
        ("Academic cybersecurity research", 1),
        ("Security tool development", 1),
        ("Threat intelligence gathering", 1),
        ("🛡️ Ethical Hacking", 0),
        ("Bug bounty preparation", 1),
        ("Red team exercises", 1),
        ("Security certification studies", 1)
    ]
    
    for use_case, level in use_cases:
        p = content.add_paragraph()
        p.text = use_case
        p.level = level
    
    # Slide 12: Future Enhancements
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    slide12.shapes.title.text = "🚀 Future Enhancements"
    
    content = slide12.placeholders[1].text_frame
    content.text = "Planned Features"
    
    future_features = [
        ("🔍 OSINT Expansion", 0),
        ("Social media intelligence integration", 1),
        ("Domain and IP intelligence tools", 1),
        ("Dark web monitoring capabilities", 1),
        ("🤖 Advanced AI Features", 0),
        ("Machine learning for threat detection", 1),
        ("Natural language processing enhancements", 1),
        ("Predictive analytics for security trends", 1),
        ("📱 Mobile Application", 0),
        ("React Native cross-platform app", 1),
        ("Offline capabilities", 1),
        ("Push notifications for real-time alerts", 1),
        ("🏢 Enterprise Features", 0),
        ("Team collaboration tools", 1),
        ("Advanced reporting and analytics", 1),
        ("Enterprise security compliance", 1)
    ]
    
    for feature, level in future_features:
        p = content.add_paragraph()
        p.text = feature
        p.level = level
    
    # Slide 13: Deployment Architecture
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    slide13.shapes.title.text = "🚀 Deployment Architecture"
    
    content = slide13.placeholders[1].text_frame
    content.text = "Production Ready Deployment"
    
    deployment = [
        ("🐳 Containerization", 0),
        ("Docker containerization for consistency", 1),
        ("Docker Compose for multi-service setup", 1),
        ("☁️ Cloud Deployment", 0),
        ("AWS/Azure/GCP compatible", 1),
        ("Auto-scaling capabilities", 1),
        ("Load balancer integration", 1),
        ("🔄 CI/CD Pipeline", 0),
        ("Automated testing and deployment", 1),
        ("Blue-green deployment strategy", 1),
        ("Rollback capabilities", 1),
        ("📊 Monitoring & Logging", 0),
        ("Application performance monitoring", 1),
        ("Error tracking and alerting", 1),
        ("Security event logging", 1)
    ]
    
    for deploy_item, level in deployment:
        p = content.add_paragraph()
        p.text = deploy_item
        p.level = level
    
    # Slide 14: Installation & Setup
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    slide14.shapes.title.text = "⚙️ Installation & Setup"
    
    content = slide14.placeholders[1].text_frame
    content.text = "Quick Start Guide"
    
    setup_steps = [
        ("1️⃣ Environment Setup", 0),
        ("Python 3.13+ installation", 1),
        ("Virtual environment creation", 1),
        ("Dependencies installation via pip", 1),
        ("2️⃣ Database Configuration", 0),
        ("Database migrations execution", 1),
        ("Sample data loading", 1),
        ("Admin user creation", 1),
        ("3️⃣ API Configuration", 0),
        ("Google Gemini API key setup", 1),
        ("Environment variables configuration", 1),
        ("4️⃣ Server Launch", 0),
        ("Development server startup", 1),
        ("Access via http://localhost:8000", 1),
        ("Admin panel at /admin/", 1)
    ]
    
    for step, level in setup_steps:
        p = content.add_paragraph()
        p.text = step
        p.level = level
    
    # Slide 15: Demo Screenshots
    slide15 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide15.shapes.title.text = "📸 Application Screenshots"
    
    # Add text box for screenshot descriptions
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5)
    
    textbox = slide15.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Key Interface Screenshots:"
    
    screenshots = [
        "🏠 Main Dashboard - Overview of dorks and statistics",
        "🔍 Dork Browser - Categorized dork exploration",
        "💬 AI Chatbot - Interactive security assistant",
        "👤 User Profile - Personal settings and API management",
        "📊 Session Dashboard - Research session management",
        "🔐 Admin Panel - Content and user management"
    ]
    
    for screenshot in screenshots:
        p = text_frame.add_paragraph()
        p.text = screenshot
        p.level = 1
    
    # Slide 16: Conclusion
    slide16 = prs.slides.add_slide(prs.slide_layouts[1])
    slide16.shapes.title.text = "🎯 Conclusion"
    
    content = slide16.placeholders[1].text_frame
    content.text = "Project Impact"
    
    conclusion_points = [
        ("✅ Successfully Delivered", 0),
        ("Comprehensive cybersecurity research platform", 1),
        ("AI-enhanced user experience", 1),
        ("Production-ready architecture", 1),
        ("Secure and scalable design", 1),
        ("🎓 Educational Value", 0),
        ("Democratizes cybersecurity knowledge", 1),
        ("Promotes ethical hacking practices", 1),
        ("Supports academic research", 1),
        ("🚀 Future Potential", 0),
        ("Extensible architecture for new features", 1),
        ("Strong foundation for enterprise expansion", 1),
        ("Community-driven development opportunities", 1),
        ("💡 Innovation", 0),
        ("Unique combination of AI and security tools", 1),
        ("User-centric design approach", 1),
        ("Modern web technologies implementation", 1)
    ]
    
    for point, level in conclusion_points:
        p = content.add_paragraph()
        p.text = point
        p.level = level
    
    # Slide 17: Thank You
    slide17 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide17.shapes.title
    subtitle = slide17.placeholders[1]
    
    title.text = "Thank You!"
    subtitle.text = "Questions & Discussion\n\n🔗 GitHub: github.com/kineticKshitij/MajorProject-V1\n📧 Contact: [Your Email]\n🌐 Demo: [Your Demo URL]"
    
    # Format thank you slide
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    """Main function to create and save the presentation"""
    print("🎯 Creating AI-Enhanced Google Dorks Toolkit Presentation...")
    
    try:
        # Create the presentation
        presentation = create_presentation()
        
        # Save the presentation
        output_file = "AI_Enhanced_Google_Dorks_Toolkit_Presentation.pptx"
        presentation.save(output_file)
        
        print(f"✅ Presentation created successfully!")
        print(f"📁 File saved as: {output_file}")
        print(f"📊 Total slides: {len(presentation.slides)}")
        
        # Get file size
        file_size = os.path.getsize(output_file) / 1024  # Size in KB
        print(f"📦 File size: {file_size:.1f} KB")
        
        print("\n🎉 Your presentation is ready!")
        print("\n📋 Presentation Contents:")
        
        slide_titles = [
            "1. Title Slide",
            "2. Project Overview",
            "3. Key Features",
            "4. Technical Architecture",
            "5. Database Architecture",
            "6. User Workflows",
            "7. Security Implementation",
            "8. AI Integration Details",
            "9. Technology Stack",
            "10. Project Metrics",
            "11. Use Cases & Applications",
            "12. Future Enhancements",
            "13. Deployment Architecture",
            "14. Installation & Setup",
            "15. Application Screenshots",
            "16. Conclusion",
            "17. Thank You & Contact"
        ]
        
        for title in slide_titles:
            print(f"   ✓ {title}")
            
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        return False
    
    return True

if __name__ == "__main__":
    main()